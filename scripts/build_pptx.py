"""PNG → PowerPoint (16:9 widescreen, 22 slides)

새 Claude 프로젝트에서 사용:
1. python scripts/render_pngs.py 먼저 실행 (PNG 생성)
2. python scripts/build_pptx.py 실행
3. ./output/heensun-marketing-strategy-2026.pptx 생성

요구 사항:
- pip install python-pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

SCRIPT_DIR = Path(__file__).resolve().parent
ROOT = SCRIPT_DIR.parent
PNG_DIR = ROOT / "png"
OUT_DIR = ROOT / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "heensun-marketing-strategy-2026.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

png_files = sorted(PNG_DIR.glob("slide-*.png"))
if not png_files:
    raise SystemExit(f"⚠ PNG 없음. 먼저 render_pngs.py 실행 필요. ({PNG_DIR})")

print(f"빌드: {len(png_files)} 슬라이드 → {OUT_PPTX}\n")
for png in png_files:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)
    print(f"  + {png.name}")

prs.save(str(OUT_PPTX))
print(f"\n완료: {OUT_PPTX}")
print(f"크기: {OUT_PPTX.stat().st_size / 1024:.1f} KB")
